from __future__ import annotations

import sys
from dataclasses import dataclass, field
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from kip.adapters.parsers.pptx import PptxParser
from kip.adapters.parsers.pptx_ocr import PptxOcrLimits
from kip.adapters.repository.memory import MemoryRepository
from kip.container import build_container
from kip.domain.models import SearchRequest
from kip.errors import ConfigurationError, ParserError
from kip.ports.ocr import OcrBlock, OcrDocument


@dataclass(slots=True)
class KoreanPictureOcr:
    name: str = "fixture-korean-ocr"
    version: str = "1"
    batches: list[tuple[Path, ...]] = field(default_factory=list)

    def recognize(self, paths: tuple[Path, ...]) -> tuple[OcrDocument, ...]:
        self.batches.append(paths)
        return tuple(
            OcrDocument(
                source_path=path,
                blocks=(
                    OcrBlock(
                        text="설비 데이터 수집 현황",
                        block_type="paragraph",
                        page=1,
                        bbox={"x": 8, "y": 16, "width": 320, "height": 44},
                        metadata={"confidence": 0.96},
                    ),
                ),
                metadata={"ocrEngine": "PP-OCRv5 korean"},
                warnings=(),
            )
            for path in paths
        )


class FailingOcr:
    name = "fixture-failing-ocr"
    version = "1"

    def recognize(self, paths: tuple[Path, ...]) -> tuple[OcrDocument, ...]:
        raise ParserError("fixture OCR unavailable")


def test_pptx_ocr_deduplicates_images_and_preserves_each_shape_locator(
    tmp_path: Path,
) -> None:
    # Given two picture shapes that reuse the same Korean screenshot bytes.
    image_path = tmp_path / "dashboard.png"
    Image.new("RGB", (800, 300), color="white").save(image_path)
    deck_path = tmp_path / "status.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    first = slide.shapes.add_picture(
        str(image_path), Inches(1), Inches(1), Inches(4), Inches(1.5)
    )
    second = slide.shapes.add_picture(
        str(image_path), Inches(1), Inches(3), Inches(4), Inches(1.5)
    )
    presentation.save(deck_path)
    ocr = KoreanPictureOcr()

    # When the structural PPTX parser enriches eligible pictures with OCR.
    extraction, units = PptxParser(ocr=ocr).parse(
        deck_path,
        artifact_id="art_pptx_ocr",
        document_id="doc_pptx_ocr",
        acl_scopes=["workspace:default"],
    )

    # Then the image is recognized once but both source shapes retain OCR evidence.
    assert len(ocr.batches) == 1
    assert len(ocr.batches[0]) == 1
    ocr_units = [unit for unit in units if unit.unit_type == "pptx_ocr"]
    assert [unit.body for unit in ocr_units] == [
        "설비 데이터 수집 현황",
        "설비 데이터 수집 현황",
    ]
    assert [unit.locator.data["shape_id"] for unit in ocr_units] == [
        first.shape_id,
        second.shape_id,
    ]
    assert ocr_units[0].locator.data["ocr_bbox_px"] == {
        "x": 8,
        "y": 16,
        "width": 320,
        "height": 44,
    }
    assert extraction.parser_name == "python-pptx+fixture-korean-ocr"
    assert extraction.metadata["ocr_unique_image_count"] == 1
    assert extraction.metadata["ocr_block_count"] == 2


def test_pptx_ocr_is_searchable_through_normal_sync(
    test_container, tmp_path: Path
) -> None:
    # Given a source deck with image-only Korean evidence and an enabled pinned OCR command.
    command = tmp_path / "fake_kordoc.py"
    command.write_text(
        """
import json
import sys
if "--version" in sys.argv:
    print("4.7.3")
else:
    for value in sys.argv[1:]:
        if value.endswith((".png", ".jpg", ".jpeg", ".webp")):
            print(json.dumps({
                "success": True,
                "fileType": "image",
                "blocks": [{
                    "type": "paragraph",
                    "text": "설비 데이터 수집 현황",
                    "pageNumber": 1,
                    "bbox": {"x": 5, "y": 7, "width": 200, "height": 30}
                }]
            }, ensure_ascii=False))
""".strip(),
        encoding="utf-8",
    )
    settings = test_container.settings
    settings.raw["sources"]["filesystem"][0]["include_extensions"].append(".pptx")
    settings.raw["parsers"]["ocr"] = {
        "timeout_seconds": 5,
        "kordoc": {
            "enabled": True,
            "argv": [
                sys.executable,
                str(command),
                "--format",
                "json",
                "--ocr",
                "--silent",
            ],
            "version_argv": [sys.executable, str(command), "--version"],
            "expected_version": "4.7.3",
        },
    }
    source_path = settings.project_root / "source" / "이미지현황.pptx"
    image_path = tmp_path / "status.png"
    Image.new("RGB", (800, 300), color="white").save(image_path)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(
        str(image_path), Inches(1), Inches(1), Inches(6), Inches(2)
    )
    presentation.save(source_path)
    source_before = source_path.stat()
    container = build_container(settings, repository=MemoryRepository())

    # When normal filesystem sync and lexical search run.
    context = container.application.operations.request_context()
    summary = container.application.ingestion.sync_filesystem(context, "fixture")
    hits = container.application.retrieval.search(
        context,
        SearchRequest(query="설비 데이터", limit=10),
    )

    # Then OCR evidence is retrievable with its source shape and the deck is unchanged.
    assert summary.inserted == 1
    assert hits
    unit = container.repository.evidence.get_content_unit(context, hits[0].unit_id)
    assert unit.unit_type == "pptx_ocr"
    assert unit.body == "설비 데이터 수집 현황"
    assert unit.locator.type == "pptx_ocr"
    source_after = source_path.stat()
    assert (source_before.st_size, source_before.st_mtime_ns) == (
        source_after.st_size,
        source_after.st_mtime_ns,
    )


def test_pptx_preserves_native_units_when_ocr_fails(tmp_path: Path) -> None:
    # Given a deck with native text and an image whose OCR command fails.
    image_path = tmp_path / "dashboard.png"
    Image.new("RGB", (800, 300), color="white").save(image_path)
    deck_path = tmp_path / "status.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(
        Inches(1), Inches(0.25), Inches(6), Inches(0.5)
    ).text = "원문 텍스트 유지"
    slide.shapes.add_picture(
        str(image_path), Inches(1), Inches(1), Inches(6), Inches(2)
    )
    presentation.save(deck_path)

    # When optional OCR enrichment fails.
    extraction, units = PptxParser(ocr=FailingOcr()).parse(
        deck_path,
        artifact_id="art_pptx_fallback",
        document_id="doc_pptx_fallback",
        acl_scopes=["workspace:default"],
    )

    # Then the structural extraction remains usable and the failure is explicit.
    assert any("원문 텍스트 유지" in unit.body for unit in units)
    assert not any(unit.unit_type == "pptx_ocr" for unit in units)
    assert extraction.status == "partial"
    assert extraction.warnings == ["OCR_FAILED: fixture OCR unavailable"]


def test_pptx_ocr_limits_reject_non_positive_values() -> None:
    # Given an image budget that cannot admit any valid work.
    # When the budget crosses the configuration boundary.
    with pytest.raises(ConfigurationError, match="max_images"):
        PptxOcrLimits(max_images=0)

    # Then invalid limits fail before a source file is opened.
