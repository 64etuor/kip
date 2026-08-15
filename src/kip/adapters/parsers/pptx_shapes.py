from __future__ import annotations

from collections.abc import Iterable
from dataclasses import dataclass
from typing import TYPE_CHECKING, Final, Protocol, assert_never, cast

from kip.domain.json_types import JsonObject, JsonValue
from kip.ids import sha256_bytes

if TYPE_CHECKING:
    from pptx.chart.chart import Chart
    from pptx.shapes.autoshape import Shape
    from pptx.shapes.base import BaseShape
    from pptx.shapes.group import GroupShape
    from pptx.shapes.picture import Picture
    from pptx.slide import Slide
    from pptx.table import Table


@dataclass(frozen=True, slots=True)
class PptxShapeRecord:
    unit_type: str
    body: str
    shape_id: int
    shape_name: str
    group_path: tuple[int, ...]
    bbox_emu: JsonObject
    z_order: tuple[int, ...]
    metadata: JsonObject


class ShapeLike(Protocol):
    @property
    def left(self) -> int | None: ...

    @property
    def top(self) -> int | None: ...

    @property
    def width(self) -> int | None: ...

    @property
    def height(self) -> int | None: ...


@dataclass(frozen=True, slots=True)
class _Transform:
    """Maps a shape's raw local (left, top, width, height) to slide EMU.

    A shape's own `<a:off>`/`<a:ext>` are always defined relative to its
    *immediate parent's* coordinate system: for a top-level shape that
    parent is the slide itself (an implicit identity transform), but for a
    shape inside a `<p:grpSp>` it is the group's *child* coordinate space
    (`chOff`/`chExt`), which PowerPoint lets an author move/resize
    independently of the group's own on-slide `off`/`ext` - e.g. moving or
    resizing a group after grouping, the most common real-world grouping
    edit, changes this mapping without touching a single child shape's raw
    XML. python-pptx's `.left`/`.top`/`.width`/`.height` return those raw,
    unconverted local values, so reading them directly for a shape nested
    in a group silently reports the wrong on-slide position/size (and
    reading order, which sorts by this same bbox) - correct only by
    coincidence, when a group has never been moved or resized after its
    child extents were first computed.
    """

    scale_x: float
    scale_y: float
    offset_x: float
    offset_y: float

    def apply(self, left: int, top: int, width: int, height: int) -> tuple[int, int, int, int]:
        return (
            round(self.offset_x + left * self.scale_x),
            round(self.offset_y + top * self.scale_y),
            round(width * self.scale_x),
            round(height * self.scale_y),
        )

    def child_transform(
        self, *, off_x: int, off_y: int, ext_cx: int, ext_cy: int, ch_off_x: int, ch_off_y: int, ch_ext_cx: int, ch_ext_cy: int
    ) -> _Transform:
        """Compute the transform for a group's direct children.

        ``off``/``ext`` are the group's own raw local position/size (in
        *this* transform's coordinate space); ``chOff``/``chExt`` are the
        group's child-coordinate-space origin/size that its children's
        ``off``/``ext`` values are expressed in. Rotation/flip on the group
        are intentionally not composed in here (consistent with how a
        shape's own rotation is already reported separately as
        ``rotation_degrees`` rather than folded into ``bbox_emu`` - see
        ``_shape_metadata``): the result is the group's un-rotated frame in
        slide-EMU terms, which keeps this a well-defined, low-risk affine
        translate+scale instead of a lossy rotated-bbox approximation.
        """
        abs_x, abs_y, abs_cx, abs_cy = self.apply(off_x, off_y, ext_cx, ext_cy)
        scale_x = (abs_cx / ch_ext_cx) if ch_ext_cx else 1.0
        scale_y = (abs_cy / ch_ext_cy) if ch_ext_cy else 1.0
        return _Transform(
            scale_x=scale_x,
            scale_y=scale_y,
            offset_x=abs_x - ch_off_x * scale_x,
            offset_y=abs_y - ch_off_y * scale_y,
        )


_IDENTITY_TRANSFORM: Final = _Transform(scale_x=1.0, scale_y=1.0, offset_x=0.0, offset_y=0.0)


def extract_shape_records(slide: Slide) -> list[PptxShapeRecord]:
    records = _extract_shapes(slide.shapes, group_path=(), z_prefix=(), transform=_IDENTITY_TRANSFORM)
    return sorted(
        records,
        key=lambda record: (
            cast(int, record.bbox_emu["top"]),
            cast(int, record.bbox_emu["left"]),
            record.z_order,
        ),
    )


def _extract_shapes(
    shapes: Iterable[BaseShape],
    *,
    group_path: tuple[int, ...],
    z_prefix: tuple[int, ...],
    transform: _Transform,
) -> list[PptxShapeRecord]:
    from pptx.shapes.autoshape import Shape
    from pptx.shapes.base import BaseShape
    from pptx.shapes.graphfrm import GraphicFrame
    from pptx.shapes.group import GroupShape
    from pptx.shapes.picture import Picture

    records: list[PptxShapeRecord] = []
    for z_order, shape in enumerate(shapes):
        bbox = _bbox(shape, transform)
        z_path = (*z_prefix, z_order)
        match shape:
            case GroupShape():
                records.extend(
                    _extract_shapes(
                        shape.shapes,
                        group_path=(*group_path, shape.shape_id),
                        z_prefix=z_path,
                        transform=_group_child_transform(shape, transform),
                    )
                )
            case GraphicFrame() if shape.has_table:
                body, table = _extract_table(shape.table)
                if body:
                    records.append(
                        PptxShapeRecord(
                            unit_type="pptx_table",
                            body=body,
                            shape_id=shape.shape_id,
                            shape_name=shape.name,
                            group_path=group_path,
                            bbox_emu=bbox,
                            z_order=z_path,
                            metadata=_shape_metadata(shape, "table", table),
                        )
                    )
            case GraphicFrame() if shape.has_chart:
                body, chart = _extract_chart(shape.chart)
                records.append(
                    PptxShapeRecord(
                        unit_type="pptx_chart",
                        body=body,
                        shape_id=shape.shape_id,
                        shape_name=shape.name,
                        group_path=group_path,
                        bbox_emu=bbox,
                        z_order=z_path,
                        metadata=_shape_metadata(shape, "chart", chart),
                    )
                )
            case Picture():
                body, image = _extract_image(shape)
                records.append(
                    PptxShapeRecord(
                        unit_type="pptx_image",
                        body=body,
                        shape_id=shape.shape_id,
                        shape_name=shape.name,
                        group_path=group_path,
                        bbox_emu=bbox,
                        z_order=z_path,
                        metadata=_shape_metadata(shape, "image", image),
                    )
                )
            case Shape() if shape.has_text_frame:
                text, paragraphs = _extract_text(shape)
                if text:
                    records.append(
                        PptxShapeRecord(
                            unit_type="pptx_text",
                            body=text,
                            shape_id=shape.shape_id,
                            shape_name=shape.name,
                            group_path=group_path,
                            bbox_emu=bbox,
                            z_order=z_path,
                            metadata=_shape_metadata(shape, "paragraphs", paragraphs),
                        )
                    )
            case BaseShape():
                continue
            case unreachable:
                assert_never(unreachable)
    return records


def _extract_table(table: Table) -> tuple[str, JsonObject]:
    rows: list[JsonValue] = []
    text_rows: list[str] = []
    for row_index, row in enumerate(table.rows):
        cells: list[JsonValue] = []
        texts: list[str] = []
        for col_index, cell in enumerate(row.cells):
            text = cell.text.strip()
            cells.append(
                {
                    "row": row_index,
                    "col": col_index,
                    "text": text,
                    "row_span": cell.span_height if cell.is_merge_origin else 1,
                    "col_span": cell.span_width if cell.is_merge_origin else 1,
                    "is_merge_origin": cell.is_merge_origin,
                    "is_spanned": cell.is_spanned,
                }
            )
            texts.append(text)
        rows.append(cells)
        text_rows.append("\t".join(texts))
    return "\n".join(text_rows).strip(), {
        "rows": len(table.rows),
        "cols": len(table.columns),
        "cells": rows,
    }


def _shape_metadata(shape: BaseShape, key: str, value: JsonValue) -> JsonObject:
    rotation = shape.rotation
    return {
        "shape_type": str(shape.shape_type),
        "rotation_degrees": float(rotation) if rotation is not None else None,
        "coordinate_space": "slide_emu",
        key: value,
    }


def _extract_chart(chart: Chart) -> tuple[str, JsonObject]:
    title = chart.chart_title.text_frame.text.strip() if chart.has_title else ""
    if len(chart.plots) == 0:
        return title or "[chart]", {"title": title, "categories": [], "series": []}
    plot_categories = chart.plots[0].categories
    categories = (
        [str(category.label) for category in plot_categories] if plot_categories is not None else []
    )
    series: list[JsonValue] = []
    series_names: list[str] = []
    series_values: list[list[float | None]] = []
    for item in chart.series:
        name = str(item.name)
        values = [float(value) if value is not None else None for value in item.values]
        series_names.append(name)
        series_values.append(values)
    point_count = max([len(categories), *(len(values) for values in series_values)])
    categories.extend("" for _ in range(point_count - len(categories)))
    for name, values in zip(series_names, series_values, strict=True):
        values.extend(None for _ in range(point_count - len(values)))
        json_values: list[JsonValue] = list(values)
        series.append({"name": name, "values": json_values})
    lines = [title] if title else []
    lines.append("Category\t" + "\t".join(series_names))
    for index, category in enumerate(categories):
        row_values = ["" if items[index] is None else str(items[index]) for items in series_values]
        lines.append("\t".join([category, *row_values]))
    category_values: list[JsonValue] = list(categories)
    return "\n".join(lines), {
        "title": title,
        "categories": category_values,
        "series": series,
    }


def _extract_image(picture: Picture) -> tuple[str, JsonObject]:
    image = picture.image
    alt_text = _alt_text(picture)
    label = alt_text or image.filename or picture.name
    return f"[image: {label}]", {
        "filename": image.filename or picture.name,
        "mime_type": image.content_type,
        "sha256": sha256_bytes(image.blob),
        "size_bytes": len(image.blob),
        "alt_text": alt_text,
    }


def _alt_text(picture: Picture) -> str:
    for element in picture.element.iter():
        if element.tag.rsplit("}", 1)[-1] != "cNvPr":
            continue
        return (element.get("descr") or element.get("title") or "").strip()
    return ""


def _extract_text(shape: Shape) -> tuple[str, list[JsonValue]]:
    paragraphs: list[JsonValue] = []
    for paragraph in shape.text_frame.paragraphs:
        runs: list[JsonValue] = []
        for run in paragraph.runs:
            runs.append(
                {
                    "text": run.text,
                    "bold": run.font.bold,
                    "italic": run.font.italic,
                    "font_size_pt": run.font.size.pt if run.font.size is not None else None,
                    "hyperlink": run.hyperlink.address,
                }
            )
        paragraphs.append({"level": paragraph.level, "text": paragraph.text, "runs": runs})
    return shape.text.strip(), paragraphs


def _bbox(shape: ShapeLike, transform: _Transform) -> JsonObject:
    left, top, width, height = transform.apply(
        _length(shape.left), _length(shape.top), _length(shape.width), _length(shape.height)
    )
    return {"left": left, "top": top, "width": width, "height": height}


def _length(value: int | None) -> int:
    return int(value) if value is not None else 0


def _group_child_transform(shape: GroupShape, transform: _Transform) -> _Transform:
    """Build the transform for the direct children of group ``shape``.

    Falls back to the group's own (already-correct) transform, unchanged,
    when the group has no explicit child-coordinate-space extent to scale
    from (an empty or malformed ``chExt``) - the same "no-op" 1:1 mapping
    `child_transform` already applies per-axis when only one axis is
    degenerate.
    """
    element = shape.element
    ch_off = element.chOff
    ch_ext = element.chExt
    return transform.child_transform(
        off_x=_length(shape.left),
        off_y=_length(shape.top),
        ext_cx=_length(shape.width),
        ext_cy=_length(shape.height),
        ch_off_x=_length(ch_off.x),
        ch_off_y=_length(ch_off.y),
        ch_ext_cx=_length(ch_ext.cx),
        ch_ext_cy=_length(ch_ext.cy),
    )
