from __future__ import annotations

import tempfile
from collections.abc import Iterable
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import TYPE_CHECKING, assert_never

from kip.adapters.parsers.pptx_units import PptxUnitContext
from kip.domain.json_types import JsonObject
from kip.domain.models import ContentUnit, EvidenceLocator
from kip.domain.text import normalize_text
from kip.errors import ConfigurationError
from kip.ids import sha256_bytes, stable_id
from kip.ports.ocr import OcrPort

if TYPE_CHECKING:
    from pptx.presentation import Presentation
    from pptx.shapes.base import BaseShape
    from pptx.shapes.picture import Picture


@dataclass(frozen=True, slots=True)
class PptxOcrLimits:
    max_images: int = 128
    max_image_bytes: int = 20 * 1024 * 1024
    max_total_bytes: int = 100 * 1024 * 1024
    min_width_px: int = 96
    min_height_px: int = 48

    def __post_init__(self) -> None:
        for field_name in (
            "max_images",
            "max_image_bytes",
            "max_total_bytes",
            "min_width_px",
            "min_height_px",
        ):
            if getattr(self, field_name) <= 0:
                raise ConfigurationError(f"invalid {field_name} OCR limit")


@dataclass(frozen=True, slots=True)
class PptxOcrOutput:
    units: tuple[ContentUnit, ...]
    warnings: tuple[str, ...]
    metadata: JsonObject


@dataclass(frozen=True, slots=True)
class PptxOcrContext:
    unit: PptxUnitContext
    limits: PptxOcrLimits
    ordinal_start: int


@dataclass(frozen=True, slots=True)
class _PictureOccurrence:
    slide: int
    slide_id: int
    shape_id: int
    group_path: tuple[int, ...]
    bbox_emu: JsonObject


@dataclass(frozen=True, slots=True)
class _PictureLocation:
    slide: int
    slide_id: int
    group_path: tuple[int, ...]


@dataclass(frozen=True, slots=True)
class _PictureCandidate:
    sha256: str
    payload: bytes
    suffix: str
    occurrences: tuple[_PictureOccurrence, ...]


def enrich_presentation(
    presentation: Presentation,
    ocr: OcrPort,
    context: PptxOcrContext,
) -> PptxOcrOutput:
    candidates, warnings, skipped = _collect_candidates(
        presentation, context.limits
    )
    if not candidates:
        return PptxOcrOutput(
            units=(),
            warnings=tuple(warnings),
            metadata={
                "ocr_unique_image_count": 0,
                "ocr_block_count": 0,
                "ocr_skipped_image_count": skipped,
            },
        )
    with tempfile.TemporaryDirectory(prefix="kip-pptx-ocr-") as temp_dir:
        paths: list[Path] = []
        candidate_by_path: dict[Path, _PictureCandidate] = {}
        for index, candidate in enumerate(candidates):
            path = Path(temp_dir) / (
                f"{index:04d}-{candidate.sha256[:16]}{candidate.suffix}"
            )
            path.write_bytes(candidate.payload)
            paths.append(path)
            candidate_by_path[path] = candidate
        documents = ocr.recognize(tuple(paths))
        units: list[ContentUnit] = []
        for document in documents:
            matched_candidate = candidate_by_path.get(document.source_path)
            if matched_candidate is None:
                continue
            warnings.extend(f"OCR_WARNING: {warning}" for warning in document.warnings)
            for occurrence in matched_candidate.occurrences:
                for block in document.blocks:
                    if not block.text.strip():
                        continue
                    ordinal = context.ordinal_start + len(units)
                    normalized = normalize_text(block.text)
                    units.append(
                        ContentUnit(
                            id=stable_id(
                                "unit", context.unit.extraction_id, str(ordinal)
                            ),
                            extraction_id=context.unit.extraction_id,
                            document_id=context.unit.document_id,
                            artifact_id=context.unit.artifact_id,
                            ordinal=ordinal,
                            unit_type="pptx_ocr",
                            title=f"slide {occurrence.slide} picture OCR",
                            body=block.text,
                            body_normalized=normalized,
                            lexical_text=normalized,
                            locator=EvidenceLocator(
                                type="pptx_ocr",
                                data={
                                    "slide": occurrence.slide,
                                    "slide_id": occurrence.slide_id,
                                    "shape_id": occurrence.shape_id,
                                    "group_path": list(occurrence.group_path),
                                    "bbox_emu": occurrence.bbox_emu,
                                    "ocr_bbox_px": block.bbox,
                                },
                            ),
                            acl_scopes=list(context.unit.acl_scopes),
                            metadata={
                                "image_sha256": matched_candidate.sha256,
                                "ocr_adapter": ocr.name,
                                "block_type": block.block_type,
                                **block.metadata,
                            },
                        )
                    )
    return PptxOcrOutput(
        units=tuple(units),
        warnings=tuple(warnings),
        metadata={
            "ocr_unique_image_count": len(candidates),
            "ocr_block_count": len(units),
            "ocr_skipped_image_count": skipped,
        },
    )


def _collect_candidates(
    presentation: Presentation,
    limits: PptxOcrLimits,
) -> tuple[list[_PictureCandidate], list[str], int]:
    payloads: dict[str, tuple[bytes, str]] = {}
    occurrences: dict[str, list[_PictureOccurrence]] = {}
    warnings: list[str] = []
    skipped = 0
    total_bytes = 0
    for slide_number, slide in enumerate(presentation.slides, start=1):
        for picture, group_path in _pictures(slide.shapes):
            width_px, height_px = picture.image.size
            if width_px < limits.min_width_px or height_px < limits.min_height_px:
                skipped += 1
                continue
            blob = picture.image.blob
            digest = sha256_bytes(blob)
            occurrence = _occurrence(
                picture,
                _PictureLocation(
                    slide=slide_number,
                    slide_id=slide.slide_id,
                    group_path=group_path,
                ),
            )
            if digest in payloads:
                occurrences[digest].append(occurrence)
                continue
            if len(blob) > limits.max_image_bytes:
                skipped += 1
                warnings.append(f"OCR_IMAGE_TOO_LARGE: shape {picture.shape_id}")
                continue
            if (
                len(payloads) >= limits.max_images
                or total_bytes + len(blob) > limits.max_total_bytes
            ):
                skipped += 1
                warnings.append(f"OCR_IMAGE_LIMIT: shape {picture.shape_id}")
                continue
            normalized = _ocr_payload(blob, picture.image.ext)
            if normalized is None:
                skipped += 1
                warnings.append(f"OCR_UNSUPPORTED_IMAGE: shape {picture.shape_id}")
                continue
            payload, suffix = normalized
            payloads[digest] = (payload, suffix)
            occurrences[digest] = [occurrence]
            total_bytes += len(blob)
    candidates = [
        _PictureCandidate(
            sha256=digest,
            payload=payload,
            suffix=suffix,
            occurrences=tuple(occurrences[digest]),
        )
        for digest, (payload, suffix) in payloads.items()
    ]
    return candidates, warnings, skipped


def _pictures(
    shapes: Iterable[BaseShape],
    group_path: tuple[int, ...] = (),
) -> Iterable[tuple[Picture, tuple[int, ...]]]:
    from pptx.shapes.base import BaseShape
    from pptx.shapes.group import GroupShape
    from pptx.shapes.picture import Picture

    for shape in shapes:
        match shape:
            case GroupShape():
                yield from _pictures(shape.shapes, (*group_path, shape.shape_id))
            case Picture():
                yield shape, group_path
            case BaseShape():
                continue
            case unreachable:
                assert_never(unreachable)


def _occurrence(
    picture: Picture,
    location: _PictureLocation,
) -> _PictureOccurrence:
    return _PictureOccurrence(
        slide=location.slide,
        slide_id=location.slide_id,
        shape_id=picture.shape_id,
        group_path=location.group_path,
        bbox_emu={
            "left": int(picture.left or 0),
            "top": int(picture.top or 0),
            "width": int(picture.width or 0),
            "height": int(picture.height or 0),
        },
    )


def _ocr_payload(blob: bytes, extension: str) -> tuple[bytes, str] | None:
    suffix = f".{extension.lower()}"
    if suffix in {".png", ".jpg", ".jpeg", ".webp"}:
        return blob, suffix
    try:
        from PIL import Image, UnidentifiedImageError
    except ImportError:
        return None
    try:
        with Image.open(BytesIO(blob)) as image:
            output = BytesIO()
            image.convert("RGB").save(output, format="PNG")
    except (OSError, UnidentifiedImageError):
        return None
    return output.getvalue(), ".png"
