from __future__ import annotations

import zipfile
from collections.abc import Mapping
from datetime import UTC, datetime
from pathlib import Path
from xml.etree import ElementTree as ET

import pytest
from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from kip.adapters.parsers.pptx import PptxParser
from kip.adapters.parsers.pptx_ooxml import scan_pptx_package
from kip.adapters.parsers.registry import ParserRegistry
from kip.domain.models import SearchRequest
from kip.errors import ParserError, ValidationError
from kip.settings import Settings


def _replace_zip_members(path: Path, replacements: Mapping[str, bytes]) -> None:
    with zipfile.ZipFile(path) as archive:
        members = {name: archive.read(name) for name in archive.namelist()}
    members.update(replacements)
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for name, content in members.items():
            archive.writestr(name, content)


def test_registry_selects_pptx_parser_for_valid_package(tmp_path: Path) -> None:
    # Given a PresentationML package with a PPTX extension.
    path = tmp_path / "status.pptx"
    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr("[Content_Types].xml", "<Types />")
        archive.writestr("ppt/presentation.xml", "<p:presentation />")
    settings = Settings(project_root=tmp_path, config_path=tmp_path / "kip.toml", raw={})

    # When the shared registry resolves the source artifact.
    parser = ParserRegistry.from_settings(settings).find(path)

    # Then the structural PPTX adapter owns the artifact.
    assert parser.name == "python-pptx"


def test_pptx_parser_preserves_text_geometry_and_notes(tmp_path: Path) -> None:
    # Given a slide with positioned Korean text and speaker notes.
    path = tmp_path / "status.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    shape.text = "프로젝트 현황"
    shape.rotation = 12.5
    notes_frame = slide.notes_slide.notes_text_frame
    assert notes_frame is not None
    notes_frame.text = "결재 대기"
    presentation.save(path)

    # When the original presentation is parsed through the adapter contract.
    extraction, units = PptxParser().parse(
        path,
        artifact_id="art_status",
        document_id="doc_status",
        acl_scopes=["workspace:default"],
    )

    # Then text and notes are separate evidence with exact slide and shape locators.
    assert extraction.status == "succeeded"
    # A clean deck (no failed parts, no replacement-character decode loss)
    # keeps the parser's historical 1.0 quality ceiling.
    assert extraction.quality_score == pytest.approx(1.0)
    assert extraction.metadata["slide_count"] == 1
    assert [unit.unit_type for unit in units] == ["pptx_text", "pptx_notes"]
    assert units[0].body == "프로젝트 현황"
    assert units[0].locator.type == "pptx_shape"
    assert units[0].locator.data == {
        "slide": 1,
        "slide_id": slide.slide_id,
        "shape_id": shape.shape_id,
        "group_path": [],
        "bbox_emu": {
            "left": Inches(1),
            "top": Inches(2),
            "width": Inches(4),
            "height": Inches(1),
        },
    }
    assert units[0].metadata["shape_type"] == "TEXT_BOX (17)"
    assert units[0].metadata["rotation_degrees"] == 12.5
    assert units[0].metadata["coordinate_space"] == "slide_emu"
    assert units[1].body == "결재 대기"
    assert units[1].locator.type == "pptx_notes"
    assert units[1].locator.data == {"slide": 1, "slide_id": slide.slide_id}


def test_pptx_parser_preserves_placeholder_text(tmp_path: Path) -> None:
    # Given a standard title-layout slide whose text lives in placeholders.
    path = tmp_path / "placeholders.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    assert slide.shapes.title is not None
    slide.shapes.title.text = "주간 프로젝트 현황"
    placeholder = slide.placeholders[1]
    placeholder.text = "품질 개선 과제"
    presentation.save(path)

    # When the presentation is parsed.
    extraction, units = PptxParser().parse(
        path,
        artifact_id="art_placeholders",
        document_id="doc_placeholders",
        acl_scopes=["workspace:default"],
    )

    # Then title and subtitle placeholders are retained as shape evidence.
    assert extraction.status == "succeeded"
    assert [unit.body for unit in units] == ["주간 프로젝트 현황", "품질 개선 과제"]
    assert all(unit.unit_type == "pptx_text" for unit in units)
    assert {unit.locator.data["shape_id"] for unit in units} == {
        slide.shapes.title.shape_id,
        placeholder.shape_id,
    }


def test_pptx_parser_preserves_merged_table_structure(tmp_path: Path) -> None:
    # Given a slide table with a horizontal merge and Korean business values.
    path = tmp_path / "schedule.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_table(2, 3, Inches(1), Inches(1), Inches(6), Inches(2))
    table = shape.table
    table.cell(0, 0).text = "일정"
    table.cell(0, 1).text = "계획"
    table.cell(0, 2).text = "실적"
    table.cell(1, 0).text = "착수"
    table.cell(1, 1).text = "2026-08-13"
    table.cell(1, 2).text = "완료"
    table.cell(0, 0).merge(table.cell(0, 1))
    presentation.save(path)

    # When the table-bearing presentation is parsed.
    _extraction, units = PptxParser().parse(
        path,
        artifact_id="art_schedule",
        document_id="doc_schedule",
        acl_scopes=["workspace:default"],
    )

    # Then the table remains a table and merge ownership is explicit.
    assert len(units) == 1
    assert units[0].unit_type == "pptx_table"
    assert units[0].body == "일정\n계획\t\t실적\n착수\t2026-08-13\t완료"
    assert units[0].locator.data["shape_id"] == shape.shape_id
    assert units[0].metadata["table"] == {
        "rows": 2,
        "cols": 3,
        "cells": [
            [
                {
                    "row": 0,
                    "col": 0,
                    "text": "일정\n계획",
                    "row_span": 1,
                    "col_span": 2,
                    "is_merge_origin": True,
                    "is_spanned": False,
                },
                {
                    "row": 0,
                    "col": 1,
                    "text": "",
                    "row_span": 1,
                    "col_span": 1,
                    "is_merge_origin": False,
                    "is_spanned": True,
                },
                {
                    "row": 0,
                    "col": 2,
                    "text": "실적",
                    "row_span": 1,
                    "col_span": 1,
                    "is_merge_origin": False,
                    "is_spanned": False,
                },
            ],
            [
                {
                    "row": 1,
                    "col": 0,
                    "text": "착수",
                    "row_span": 1,
                    "col_span": 1,
                    "is_merge_origin": False,
                    "is_spanned": False,
                },
                {
                    "row": 1,
                    "col": 1,
                    "text": "2026-08-13",
                    "row_span": 1,
                    "col_span": 1,
                    "is_merge_origin": False,
                    "is_spanned": False,
                },
                {
                    "row": 1,
                    "col": 2,
                    "text": "완료",
                    "row_span": 1,
                    "col_span": 1,
                    "is_merge_origin": False,
                    "is_spanned": False,
                },
            ],
        ],
    }


def test_pptx_parser_extracts_charts_images_groups_and_links(tmp_path: Path) -> None:
    # Given a slide mixing a data chart, described image, grouped text, and external link.
    image_path = tmp_path / "process.png"
    Image.new("RGB", (8, 8), color="white").save(image_path)
    path = tmp_path / "mixed.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = ["1Q", "2Q"]
    chart_data.add_series("계획", (10, 12))
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(3),
        chart_data,
    )
    chart_shape.chart.has_title = True
    chart_shape.chart.chart_title.text_frame.text = "월간 생산량"
    picture = slide.shapes.add_picture(str(image_path), Inches(6), Inches(1), Inches(1), Inches(1))
    picture.element._nvXxPr.cNvPr.set("descr", "공정 흐름")
    group = slide.shapes.add_group_shape()
    grouped_text = group.shapes.add_textbox(Inches(1), Inches(5), Inches(3), Inches(1))
    grouped_text.text = "그룹 내부 상태"
    linked_text = slide.shapes.add_textbox(Inches(5), Inches(5), Inches(3), Inches(1))
    paragraph = linked_text.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = "상세 현황"
    run.font.bold = True
    run.hyperlink.address = "https://example.com/status"
    presentation.save(path)

    # When the mixed-content slide is parsed.
    extraction, units = PptxParser().parse(
        path,
        artifact_id="art_mixed",
        document_id="doc_mixed",
        acl_scopes=["workspace:default"],
    )

    # Then every meaningful shape is retained without fetching the external link.
    assert extraction.warnings == []
    assert [unit.unit_type for unit in units] == [
        "pptx_chart",
        "pptx_image",
        "pptx_text",
        "pptx_text",
    ]
    assert units[0].metadata["chart"] == {
        "title": "월간 생산량",
        "categories": ["1Q", "2Q"],
        "series": [{"name": "계획", "values": [10.0, 12.0]}],
    }
    assert units[1].body == "[image: 공정 흐름]"
    assert units[1].metadata["image"]["filename"] == "image.png"
    assert units[1].metadata["image"]["mime_type"] == "image/png"
    assert len(units[1].metadata["image"]["sha256"]) == 64
    assert units[2].locator.data["group_path"] == [group.shape_id]
    assert units[2].locator.data["shape_id"] == grouped_text.shape_id
    assert units[3].metadata["paragraphs"] == [
        {
            "level": 0,
            "text": "상세 현황",
            "runs": [
                {
                    "text": "상세 현황",
                    "bold": True,
                    "italic": None,
                    "font_size_pt": None,
                    "hyperlink": "https://example.com/status",
                }
            ],
        }
    ]


def test_pptx_parser_converts_grouped_shape_geometry_to_slide_absolute_coordinates(
    tmp_path: Path,
) -> None:
    # Given a text box grouped, then the *group itself* moved and resized
    # after grouping - the ordinary "select shapes, group, then drag/resize
    # the group" edit. The child text box's own <a:off>/<a:ext> in the XML
    # stay exactly where they were before grouping (1in, 1in, 2in, 1in);
    # only the group's chOff/chExt -> off/ext mapping changes to reflect
    # where the group (and therefore its content) now actually sits.
    path = tmp_path / "grouped.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    child = group.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    child.text = "이동된 그룹 내부"
    group.left = Inches(5)
    group.top = Inches(5)
    group.width = Inches(4)
    group.height = Inches(2)
    presentation.save(path)

    # When the presentation is parsed.
    _extraction, units = PptxParser().parse(
        path,
        artifact_id="art_grouped",
        document_id="doc_grouped",
        acl_scopes=["workspace:default"],
    )

    # Then the child's reported bbox reflects its true on-slide position
    # (the group's new frame, since the child fills the whole group here) -
    # not its raw, pre-move local offset inside the group. Previously this
    # read back as {left: 1in, top: 1in, width: 2in, height: 1in}: the raw
    # group-local XML values, silently wrong for any group ever moved or
    # resized after grouping (and mislabeled "coordinate_space": "slide_emu"
    # regardless).
    assert len(units) == 1
    assert units[0].locator.data["bbox_emu"] == {
        "left": Inches(5),
        "top": Inches(5),
        "width": Inches(4),
        "height": Inches(2),
    }
    assert units[0].metadata["coordinate_space"] == "slide_emu"


def test_pptx_parser_reports_untranscribed_embedded_media(tmp_path: Path) -> None:
    # Given a slide with one embedded video and its poster frame.
    movie_path = tmp_path / "process.mp4"
    movie_path.write_bytes(b"not-a-real-video")
    poster_path = tmp_path / "poster.png"
    Image.new("RGB", (8, 8), color="white").save(poster_path)
    path = tmp_path / "media.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_movie(
        str(movie_path),
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(3),
        poster_frame_image=str(poster_path),
        mime_type="video/mp4",
    )
    presentation.save(path)

    # When the presentation is parsed without a media transcription adapter.
    extraction, units = PptxParser().parse(
        path,
        artifact_id="art_media",
        document_id="doc_media",
        acl_scopes=["workspace:default"],
    )

    # Then the omitted media is counted once and the extraction is visibly partial.
    assert units == []
    assert extraction.status == "partial"
    assert extraction.metadata["media_object_count"] == 1
    assert extraction.warnings == ["SKIPPED_MEDIA: 1 media objects not transcribed"]


def test_pptx_parser_aligns_sparse_chart_series_without_dropping_categories(
    tmp_path: Path,
) -> None:
    # Given a chart cache whose series has fewer points than its category axis.
    path = tmp_path / "sparse-chart.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = ["1Q", "2Q", "3Q"]
    chart_data.add_series("실적", (10, 12))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(3),
        chart_data,
    )
    presentation.save(path)

    # When the sparse chart is parsed.
    extraction, units = PptxParser().parse(
        path,
        artifact_id="art_sparse_chart",
        document_id="doc_sparse_chart",
        acl_scopes=["workspace:default"],
    )

    # Then the missing point is an explicit null and the remaining file succeeds.
    assert extraction.status == "succeeded"
    assert units[0].metadata["chart"] == {
        "title": "",
        "categories": ["1Q", "2Q", "3Q"],
        "series": [{"name": "실적", "values": [10.0, 12.0, None]}],
    }
    assert units[0].body == "Category\t실적\n1Q\t10.0\n2Q\t12.0\n3Q\t"


def test_pptx_parser_keeps_chart_shape_when_plot_cache_is_absent(tmp_path: Path) -> None:
    # Given a chart shape whose OOXML plot area has no plot cache.
    path = tmp_path / "empty-chart.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = ["1Q"]
    chart_data.add_series("계획", (10,))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(3),
        chart_data,
    )
    presentation.save(path)
    with zipfile.ZipFile(path) as archive:
        chart_xml = ET.fromstring(archive.read("ppt/charts/chart1.xml"))
    plot_area = next(element for element in chart_xml.iter() if element.tag.endswith("plotArea"))
    for element in list(plot_area):
        if element.tag.endswith("Chart"):
            plot_area.remove(element)
    _replace_zip_members(path, {"ppt/charts/chart1.xml": ET.tostring(chart_xml)})

    # When the presentation is parsed.
    extraction, units = PptxParser().parse(
        path,
        artifact_id="art_empty_chart",
        document_id="doc_empty_chart",
        acl_scopes=["workspace:default"],
    )

    # Then the source shape remains explicit instead of aborting the file.
    assert extraction.status == "succeeded"
    assert units[0].unit_type == "pptx_chart"
    assert units[0].body == "[chart]"
    assert units[0].metadata["chart"] == {"title": "", "categories": [], "series": []}


def test_pptx_parser_keeps_z_order_separate_from_reading_order(tmp_path: Path) -> None:
    # Given shapes added back-to-front in an order different from their visual reading order.
    path = tmp_path / "order.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    lower = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(3), Inches(1))
    lower.text = "두 번째"
    upper = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    upper.text = "첫 번째"
    presentation.save(path)

    # When the slide is parsed for retrieval order.
    _extraction, units = PptxParser().parse(
        path,
        artifact_id="art_order",
        document_id="doc_order",
        acl_scopes=["workspace:default"],
    )

    # Then evidence follows geometry while the source z-order remains independently recorded.
    assert [unit.body for unit in units] == ["첫 번째", "두 번째"]
    assert [unit.metadata["reading_order"] for unit in units] == [0, 1]
    assert [unit.metadata["z_order"] for unit in units] == [[1], [0]]


def test_pptx_parser_extracts_hidden_slide_comments_and_smartart(tmp_path: Path) -> None:
    # Given a hidden slide with a legacy comment and SmartArt diagram text.
    path = tmp_path / "review.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    shape.text = "검토 현황"
    presentation.save(path)
    with zipfile.ZipFile(path) as archive:
        slide_xml = archive.read("ppt/slides/slide1.xml").replace(
            b"<p:sld ", b'<p:sld show="0" ', 1
        )
        rels_xml = archive.read("ppt/slides/_rels/slide1.xml.rels").replace(
            b"</Relationships>",
            (
                b'<Relationship Id="rIdComments" '
                b'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments" '
                b'Target="../comments/comment1.xml"/>'
                b'<Relationship Id="rIdDiagram" '
                b'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramData" '
                b'Target="../diagrams/data1.xml"/>'
                b"</Relationships>"
            ),
        )
    _replace_zip_members(
        path,
        {
            "ppt/slides/slide1.xml": slide_xml,
            "ppt/slides/_rels/slide1.xml.rels": rels_xml,
            "ppt/comments/comment1.xml": (
                '<p:cmLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
                '<p:cm authorId="0" idx="1" dt="2026-08-13T09:30:00Z">'
                '<p:pos x="0" y="0"/><p:text>승인 필요</p:text>'
                "</p:cm></p:cmLst>"
            ).encode(),
            "ppt/commentAuthors.xml": (
                '<p:cmAuthorLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
                '<p:cmAuthor id="0" name="홍길동" initials="HG" lastIdx="1" clrIdx="0"/>'
                "</p:cmAuthorLst>"
            ).encode(),
            "ppt/diagrams/data1.xml": (
                '<dgm:dataModel xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram">'
                "<dgm:ptLst><dgm:pt><dgm:t>검사 대기</dgm:t></dgm:pt></dgm:ptLst>"
                "</dgm:dataModel>"
            ).encode(),
        },
    )

    # When the OOXML-enriched presentation is parsed.
    extraction, units = PptxParser().parse(
        path,
        artifact_id="art_review",
        document_id="doc_review",
        acl_scopes=["workspace:default"],
    )

    # Then hidden state, attributed comment, and diagram evidence are explicit.
    assert extraction.metadata["hidden_slide_count"] == 1
    assert [unit.unit_type for unit in units] == [
        "pptx_text",
        "pptx_comment",
        "pptx_diagram",
    ]
    assert units[0].metadata["hidden_slide"] is True
    assert units[1].body == "홍길동: 승인 필요"
    assert units[1].locator.data == {"slide": 1, "comment_index": 1}
    assert units[1].metadata["author"] == "홍길동"
    assert units[1].metadata["created_at"] == "2026-08-13T09:30:00Z"
    assert units[2].body == "검사 대기"
    assert units[2].locator.data == {"slide": 1, "part": "ppt/diagrams/data1.xml"}


def test_pptx_parser_preserves_json_safe_document_properties(tmp_path: Path) -> None:
    # Given a presentation with business metadata and timezone-aware timestamps.
    path = tmp_path / "properties.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    properties = presentation.core_properties
    properties.title = "솔라엣지 프로젝트 현황"
    properties.author = "품질혁신팀"
    properties.subject = "주간 보고"
    properties.keywords = "일정,품질"
    properties.category = "내부"
    properties.created = datetime(2026, 8, 13, 9, 30, tzinfo=UTC)
    properties.modified = datetime(2026, 8, 13, 10, 45, tzinfo=UTC)
    presentation.save(path)

    # When package properties cross the public extraction boundary.
    extraction, _units = PptxParser().parse(
        path,
        artifact_id="art_properties",
        document_id="doc_properties",
        acl_scopes=["workspace:default"],
    )

    # Then metadata remains meaningful and strictly JSON-safe.
    assert extraction.metadata["document_properties"] == {
        "title": "솔라엣지 프로젝트 현황",
        "author": "품질혁신팀",
        "subject": "주간 보고",
        "keywords": "일정,품질",
        "category": "내부",
        "created": "2026-08-13T09:30:00+00:00",
        "modified": "2026-08-13T10:45:00+00:00",
    }


def test_pptx_parser_keeps_primary_content_when_comment_part_is_missing(tmp_path: Path) -> None:
    # Given a valid slide whose optional comment relationship points to a missing part.
    path = tmp_path / "missing-comment.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = "본문 유지"
    presentation.save(path)
    with zipfile.ZipFile(path) as archive:
        rels_xml = archive.read("ppt/slides/_rels/slide1.xml.rels").replace(
            b"</Relationships>",
            (
                b'<Relationship Id="rIdMissingComment" '
                b'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments" '
                b'Target="../comments/missing.xml"/>'
                b"</Relationships>"
            ),
        )
    _replace_zip_members(path, {"ppt/slides/_rels/slide1.xml.rels": rels_xml})

    # When the presentation is parsed.
    extraction, units = PptxParser().parse(
        path,
        artifact_id="art_partial",
        document_id="doc_partial",
        acl_scopes=["workspace:default"],
    )

    # Then primary evidence survives and the optional-part failure is explicit.
    assert extraction.status == "partial"
    assert [unit.body for unit in units] == ["본문 유지"]
    assert any(
        warning.startswith("PARTIAL_PARSE slide 1 comments:") for warning in extraction.warnings
    )
    # And quality is scaled down by the failed-part fraction (the slide
    # itself plus its one attempted comments part: 1 of 2 parts failed),
    # below the clean-deck 1.0 ceiling.
    assert extraction.quality_score == pytest.approx(0.5)


def test_scan_pptx_package_counts_processed_and_failed_parts(tmp_path: Path) -> None:
    # Given a single-slide deck whose only comments relationship is broken.
    path = tmp_path / "missing-comment.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = "본문 유지"
    presentation.save(path)
    with zipfile.ZipFile(path) as archive:
        rels_xml = archive.read("ppt/slides/_rels/slide1.xml.rels").replace(
            b"</Relationships>",
            (
                b'<Relationship Id="rIdMissingComment" '
                b'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments" '
                b'Target="../comments/missing.xml"/>'
                b"</Relationships>"
            ),
        )
    _replace_zip_members(path, {"ppt/slides/_rels/slide1.xml.rels": rels_xml})

    # When the package is scanned.
    package_info = scan_pptx_package(path)

    # Then both the slide part and the failed comments part are counted.
    assert package_info.processed_part_count == 2
    assert package_info.failed_part_count == 1


def test_scan_pptx_package_reports_no_failed_parts_for_a_clean_deck(tmp_path: Path) -> None:
    path = tmp_path / "clean.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(path)

    package_info = scan_pptx_package(path)

    assert package_info.processed_part_count == 1
    assert package_info.failed_part_count == 0


def test_pptx_parser_penalizes_replacement_characters_in_extracted_text(
    tmp_path: Path,
) -> None:
    # Given a slide whose extracted text already contains Unicode replacement
    # characters (as a prior encoding failure would leave behind).
    path = tmp_path / "garbled.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = "����"
    presentation.save(path)

    # When the presentation is parsed.
    extraction, units = PptxParser().parse(
        path,
        artifact_id="art_garbled",
        document_id="doc_garbled",
        acl_scopes=["workspace:default"],
    )

    # Then no part failed, but quality is reduced below the clean 1.0 ceiling
    # by the replacement-character ratio in the extracted text.
    assert extraction.status == "succeeded"
    assert units
    assert extraction.quality_score < 1.0


def test_pptx_parser_rejects_archive_over_expansion_limit(tmp_path: Path) -> None:
    # Given a highly compressible PresentationML-shaped archive.
    path = tmp_path / "oversized.pptx"
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("ppt/presentation.xml", b"0" * 1_000_000)

    # When the parser inspects the package.
    # Then it stops before expanding the archive into application objects.
    with pytest.raises(ValidationError, match="decompression limits"):
        PptxParser().parse(
            path,
            artifact_id="art_oversized",
            document_id="doc_oversized",
            acl_scopes=["workspace:default"],
        )


def test_registry_routes_corrupt_pptx_to_parser_with_clear_error(tmp_path: Path) -> None:
    # Given a file with a .pptx extension whose content is not a valid ZIP
    # archive at all (supports() previously opened the archive and
    # swallowed BadZipFile into a bare False, making the registry report
    # the misleading "no parser registered for .pptx").
    path = tmp_path / "corrupt.pptx"
    path.write_bytes(b"not a zip archive")
    settings = Settings(project_root=tmp_path, config_path=tmp_path / "kip.toml", raw={})

    # When the registry resolves a parser purely from the extension.
    parser = ParserRegistry.from_settings(settings).find(path)

    # Then the structural PPTX adapter still claims the file...
    assert parser.name == "python-pptx"
    # ...and parsing it raises the clear, typed PPTX error instead of the
    # registry ever seeing a "no parser registered" failure.
    with pytest.raises(ParserError, match="PPTX parse failed"):
        parser.parse(
            path,
            artifact_id="art_corrupt",
            document_id="doc_corrupt",
            acl_scopes=["workspace:default"],
        )


def test_pptx_parser_returns_typed_error_for_malformed_core_xml(tmp_path: Path) -> None:
    # Given a PresentationML package with malformed core presentation XML.
    path = tmp_path / "malformed.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(path)
    _replace_zip_members(path, {"ppt/presentation.xml": b"<p:presentation"})

    # When the parser reaches the malformed package boundary.
    # Then callers receive the stable parser error instead of an XML implementation error.
    with pytest.raises(ParserError, match="PPTX parse failed"):
        PptxParser().parse(
            path,
            artifact_id="art_malformed",
            document_id="doc_malformed",
            acl_scopes=["workspace:default"],
        )


def test_pptx_parser_returns_typed_error_for_malformed_slide_xml(tmp_path: Path) -> None:
    # Given a package with a valid manifest but malformed primary slide XML.
    path = tmp_path / "malformed-slide.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(path)
    _replace_zip_members(path, {"ppt/slides/slide1.xml": b"<p:sld"})

    # When python-pptx opens the corrupt primary evidence part.
    # Then its implementation-specific XML exception is translated at the adapter boundary.
    with pytest.raises(ParserError, match="PPTX parse failed"):
        PptxParser().parse(
            path,
            artifact_id="art_malformed_slide",
            document_id="doc_malformed_slide",
            acl_scopes=["workspace:default"],
        )


def test_pptx_sync_and_search_use_structured_units(test_container) -> None:
    # Given a configured filesystem source containing a Korean status presentation.
    test_container.settings.raw["sources"]["filesystem"][0]["include_extensions"].append(".pptx")
    path = test_container.settings.project_root / "source" / "주간현황.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = "검사 승인 대기"
    presentation.save(path)

    # When the normal ingestion and retrieval services process the source.
    context = test_container.application.operations.request_context()
    summary = test_container.application.ingestion.sync_filesystem(context, "fixture")
    hits = test_container.application.retrieval.search(
        context,
        SearchRequest(query="검사 승인", limit=10),
    )

    # Then the presentation is indexed through the shared parser registry as exact shape evidence.
    assert summary.inserted == 1
    assert hits
    unit = test_container.repository.evidence.get_content_unit(context, hits[0].unit_id)
    assert unit.unit_type == "pptx_text"
    assert unit.body == "검사 승인 대기"
    assert unit.locator.type == "pptx_shape"
