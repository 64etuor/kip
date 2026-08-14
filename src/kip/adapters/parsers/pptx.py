from __future__ import annotations

import zipfile
from pathlib import Path
from typing import Final
from xml.etree import ElementTree as ET

from kip.adapters.parsers.pptx_ocr import (
    PptxOcrContext,
    PptxOcrLimits,
    enrich_presentation,
)
from kip.adapters.parsers.pptx_ooxml import scan_pptx_package
from kip.adapters.parsers.pptx_units import (
    PptxSlideContext,
    PptxUnitContext,
    notes_unit,
    shape_units,
    supplemental_units,
)
from kip.domain.models import ContentUnit, ExtractionRun
from kip.errors import DependencyUnavailableError, ParserError
from kip.ids import new_id, sha256_bytes
from kip.ports.ocr import OcrPort

_PPTX_EXTENSIONS: Final = frozenset({".pptx", ".pptm", ".ppsx", ".ppsm", ".potx", ".potm"})
_DEFAULT_OCR_LIMITS: Final = PptxOcrLimits()


class PptxParser:
    name = "python-pptx"
    version = "1.0"

    def __init__(
        self,
        ocr: OcrPort | None = None,
        ocr_limits: PptxOcrLimits = _DEFAULT_OCR_LIMITS,
    ) -> None:
        self._ocr = ocr
        self._ocr_limits = ocr_limits

    def supports(self, path: Path) -> bool:
        if path.suffix.lower() not in _PPTX_EXTENSIONS:
            return False
        try:
            with zipfile.ZipFile(path) as archive:
                return "ppt/presentation.xml" in archive.namelist()
        except (OSError, zipfile.BadZipFile):
            return False

    def parse(
        self,
        path: Path,
        *,
        artifact_id: str,
        document_id: str,
        acl_scopes: list[str],
    ) -> tuple[ExtractionRun, list[ContentUnit]]:
        try:
            from pptx import Presentation
            from pptx.exc import PythonPptxError
        except ImportError as exc:
            raise DependencyUnavailableError(
                "Install the extractors extra for PPTX parsing"
            ) from exc

        try:
            package_info = scan_pptx_package(path)
            presentation = Presentation(str(path))
        except (
            OSError,
            ValueError,
            KeyError,
            ET.ParseError,
            SyntaxError,
            PythonPptxError,
            zipfile.BadZipFile,
        ) as exc:
            raise ParserError(f"PPTX parse failed: {path}: {exc}") from exc

        extraction_id = new_id("ext")
        context = PptxUnitContext(
            extraction_id=extraction_id,
            artifact_id=artifact_id,
            document_id=document_id,
            acl_scopes=tuple(acl_scopes),
            path=path,
        )
        units: list[ContentUnit] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            hidden = slide_number in package_info.hidden_slides
            slide_context = PptxSlideContext(
                unit=context,
                slide=slide,
                slide_number=slide_number,
                ordinal_start=len(units),
                hidden=hidden,
            )
            units.extend(shape_units(slide_context))
            slide_context = PptxSlideContext(
                unit=context,
                slide=slide,
                slide_number=slide_number,
                ordinal_start=len(units),
                hidden=hidden,
            )
            notes = notes_unit(slide_context)
            if notes is not None:
                units.append(notes)
            slide_context = PptxSlideContext(
                unit=context,
                slide=slide,
                slide_number=slide_number,
                ordinal_start=len(units),
                hidden=hidden,
            )
            units.extend(supplemental_units(package_info, slide_context))

        warnings = list(package_info.warnings)
        if package_info.embedded_object_count:
            warnings.append(
                f"SKIPPED_OLE: {package_info.embedded_object_count} embedded objects not expanded"
            )
        if package_info.media_object_count:
            warnings.append(
                f"SKIPPED_MEDIA: {package_info.media_object_count} media objects not transcribed"
            )
        parser_name = self.name
        ocr_metadata = {}
        if self._ocr is not None:
            parser_name = f"{self.name}+{self._ocr.name}"
            try:
                ocr_output = enrich_presentation(
                    presentation,
                    self._ocr,
                    PptxOcrContext(
                        unit=context,
                        limits=self._ocr_limits,
                        ordinal_start=len(units),
                    ),
                )
            except ParserError as exc:
                warnings.append(f"OCR_FAILED: {exc}")
            else:
                units.extend(ocr_output.units)
                warnings.extend(ocr_output.warnings)
                ocr_metadata = ocr_output.metadata
        body = "\n".join(unit.body for unit in units)
        return (
            ExtractionRun(
                id=extraction_id,
                artifact_id=artifact_id,
                parser_name=parser_name,
                parser_version=self.version,
                status="partial" if warnings else "succeeded",
                quality_score=1.0 if units else 0.0,
                output_hash=sha256_bytes(body.encode("utf-8")),
                warnings=warnings,
                metadata={
                    "slide_count": len(presentation.slides),
                    "hidden_slide_count": len(package_info.hidden_slides),
                    "animation_slide_count": package_info.animation_slide_count,
                    "transition_slide_count": package_info.transition_slide_count,
                    "embedded_object_count": package_info.embedded_object_count,
                    "media_object_count": package_info.media_object_count,
                    "external_relationship_count": package_info.external_relationship_count,
                    "document_properties": package_info.document_properties,
                    **ocr_metadata,
                },
            ),
            units,
        )
