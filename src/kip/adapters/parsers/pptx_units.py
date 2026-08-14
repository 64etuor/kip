from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import TYPE_CHECKING

from kip.adapters.parsers.pptx_ooxml import PptxPackageInfo
from kip.adapters.parsers.pptx_shapes import extract_shape_records
from kip.domain.json_types import JsonObject
from kip.domain.models import ContentUnit, EvidenceLocator
from kip.domain.text import normalize_text
from kip.ids import stable_id

if TYPE_CHECKING:
    from pptx.slide import Slide


@dataclass(frozen=True, slots=True)
class PptxUnitContext:
    extraction_id: str
    artifact_id: str
    document_id: str
    acl_scopes: tuple[str, ...]
    path: Path


@dataclass(frozen=True, slots=True)
class PptxSlideContext:
    unit: PptxUnitContext
    slide: Slide
    slide_number: int
    ordinal_start: int
    hidden: bool


@dataclass(frozen=True, slots=True)
class _SupplementalSpec:
    ordinal: int
    unit_type: str
    title_suffix: str
    body: str
    locator_type: str
    locator_data: JsonObject
    metadata: JsonObject


def shape_units(context: PptxSlideContext) -> list[ContentUnit]:
    units: list[ContentUnit] = []
    for reading_order, record in enumerate(extract_shape_records(context.slide)):
        ordinal = context.ordinal_start + len(units)
        normalized = normalize_text(record.body)
        metadata = dict(record.metadata)
        metadata.update(
            {
                "shape_name": record.shape_name,
                "z_order": list(record.z_order),
                "reading_order": reading_order,
                "hidden_slide": context.hidden,
            }
        )
        units.append(
            ContentUnit(
                id=stable_id("unit", context.unit.extraction_id, str(ordinal)),
                extraction_id=context.unit.extraction_id,
                document_id=context.unit.document_id,
                artifact_id=context.unit.artifact_id,
                ordinal=ordinal,
                unit_type=record.unit_type,
                title=f"{context.unit.path.name} - slide {context.slide_number}",
                body=record.body,
                body_normalized=normalized,
                lexical_text=normalized,
                locator=EvidenceLocator(
                    type="pptx_shape",
                    data={
                        "slide": context.slide_number,
                        "slide_id": context.slide.slide_id,
                        "shape_id": record.shape_id,
                        "group_path": list(record.group_path),
                        "bbox_emu": record.bbox_emu,
                    },
                ),
                acl_scopes=list(context.unit.acl_scopes),
                metadata=metadata,
            )
        )
    return units


def notes_unit(context: PptxSlideContext) -> ContentUnit | None:
    notes = _notes_text(context.slide)
    if not notes:
        return None
    normalized = normalize_text(notes)
    return ContentUnit(
        id=stable_id("unit", context.unit.extraction_id, str(context.ordinal_start)),
        extraction_id=context.unit.extraction_id,
        document_id=context.unit.document_id,
        artifact_id=context.unit.artifact_id,
        ordinal=context.ordinal_start,
        unit_type="pptx_notes",
        title=f"{context.unit.path.name} - slide {context.slide_number} notes",
        body=notes,
        body_normalized=normalized,
        lexical_text=normalized,
        locator=EvidenceLocator(
            type="pptx_notes",
            data={"slide": context.slide_number, "slide_id": context.slide.slide_id},
        ),
        acl_scopes=list(context.unit.acl_scopes),
        metadata={"hidden_slide": context.hidden},
    )


def supplemental_units(
    package_info: PptxPackageInfo,
    context: PptxSlideContext,
) -> list[ContentUnit]:
    units: list[ContentUnit] = []
    for comment in package_info.comments:
        if comment.slide != context.slide_number:
            continue
        body = f"{comment.author}: {comment.text}" if comment.author else comment.text
        units.append(
            _supplemental_unit(
                context,
                _SupplementalSpec(
                    ordinal=context.ordinal_start + len(units),
                    unit_type="pptx_comment",
                    title_suffix="comment",
                    body=body,
                    locator_type="pptx_comment",
                    locator_data={
                        "slide": context.slide_number,
                        "comment_index": comment.index,
                    },
                    metadata={"author": comment.author, "created_at": comment.created_at},
                ),
            )
        )
    for diagram in package_info.diagrams:
        if diagram.slide != context.slide_number:
            continue
        units.append(
            _supplemental_unit(
                context,
                _SupplementalSpec(
                    ordinal=context.ordinal_start + len(units),
                    unit_type="pptx_diagram",
                    title_suffix="diagram",
                    body=diagram.text,
                    locator_type="pptx_part",
                    locator_data={"slide": context.slide_number, "part": diagram.part},
                    metadata={},
                ),
            )
        )
    return units


def _supplemental_unit(
    context: PptxSlideContext,
    spec: _SupplementalSpec,
) -> ContentUnit:
    normalized = normalize_text(spec.body)
    return ContentUnit(
        id=stable_id("unit", context.unit.extraction_id, str(spec.ordinal)),
        extraction_id=context.unit.extraction_id,
        document_id=context.unit.document_id,
        artifact_id=context.unit.artifact_id,
        ordinal=spec.ordinal,
        unit_type=spec.unit_type,
        title=(f"{context.unit.path.name} - slide {context.slide_number} {spec.title_suffix}"),
        body=spec.body,
        body_normalized=normalized,
        lexical_text=normalized,
        locator=EvidenceLocator(type=spec.locator_type, data=spec.locator_data),
        acl_scopes=list(context.unit.acl_scopes),
        metadata=spec.metadata,
    )


def _notes_text(slide: Slide) -> str:
    if not slide.has_notes_slide:
        return ""
    text_frame = slide.notes_slide.notes_text_frame
    return text_frame.text.strip() if text_frame is not None else ""
